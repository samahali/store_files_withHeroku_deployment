from django.shortcuts import render
from django.http import JsonResponse
from .models import Document
from pptx import Presentation
import os
from django.conf import settings
from PyPDF2 import PdfFileReader

class FilesViews:
    def files_management(request):
        return render(request, template_name="store/upload_files.html")

    def search_in_files(request):
        try:
            if request.method == 'POST':
                searched_data =request.POST['search_file'].lower()
                documents = Document.objects.all()
                matched_documents=[]
                for document in documents:
                    if searched_data in document.document.name.split('/')[1].lower():
                        matched_documents +=[document]
                        continue
                    pdf_path = os.path.join(settings.MEDIA_ROOT, document.document.name)
                    if pdf_path.endswith('.pdf'):
                        search_in_file = FilesViews.search_in_pdf_file(pdf_path, searched_data)
                        if search_in_file:
                            matched_documents += [document]
                            continue
                    elif pdf_path.endswith('.pptx'):
                        search_in_file = FilesViews.search_in_pptx_file(pdf_path, searched_data)
                        if search_in_file:
                            matched_documents += [document]
                            continue
                if matched_documents:
                    return render(request, template_name="store/matched_documents.html", context={'documents': matched_documents})
                else:
                    return JsonResponse({'error':True,'message':'No items match your search.'})
        except Exception as e:
            return JsonResponse({'error':True,'message':'an application error has occurred, please contact your administrator.'})
    @staticmethod
    def search_in_pdf_file(pdf_path,searched_data):
        pdfFileObj = open(pdf_path, 'rb')
        pdfReader = PdfFileReader(pdfFileObj)
        num_pages = pdfReader.numPages
        count = 0
        text = ""  # The while loop will read each page.
        while count < num_pages:
            pageObj = pdfReader.getPage(count)
            count += 1
            text += pageObj.extractText()
            if searched_data in text.lower():
                return True
        return False
    @staticmethod
    def search_in_pptx_file(pdf_path,searched_data):
        prs = Presentation(pdf_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and searched_data in shape.text.lower():
                    return True
        return False
    def add_file(request):
        try:
            if request.method == 'POST':
                message = {"message":"Document Was Saved Successfully"}
                document_file = request.FILES['uploadedFile']
                document_file_extension = document_file.name
                if document_file_extension.endswith('.pptx') or document_file_extension.endswith('.pdf'):
                    Document.objects.create(document=document_file)
                else:
                    message["message"] = "Sorry, Files with extension .pdf,pptx only Allowed to Upload"
                    message['error'] = True
                return JsonResponse(message)
        except Exception as e:
            return JsonResponse({'error':True,'message':'an application error has occurred, please contact your administrator.'})
